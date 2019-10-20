from win32com.client import Dispatch, constants
from os.path import realpath
from zipfile import ZipFile
from lxml import etree
from pptx import Presentation

import sys


def pptx_slide_get_id_list(sldRoot):

    ret = []

    # Get slide id list
    sldIdLst = sldRoot.findall('.//{*}sldId')
    for i in range(len(sldIdLst)):
        ret.append(sldIdLst[i].get('id'))

    return ret


def pptx_slide_get_section_info(sectNode):

    # Get section name and child slide id list
    return {
        'name': sectNode.get('name'),
        'list': pptx_slide_get_id_list(sectNode)
    }


def pptx_slide_get_list(pptXml):

    sldList = []

    # Find section node
    ret = pptXml.findall('.//{*}section')
    if len(ret) > 0:
        for child in ret:
            sldList.append(pptx_slide_get_section_info(child))

    else:
        sldList.append({
            'name': None,
            'list': pptx_slide_get_id_list(pptXml)
        })

    return sldList


def pptx_slide_get_id_index_map(pptXml):

    sldIdMap = {}

    # Get slide list node
    sldRoot = None
    ret = pptXml.xpath('/p:presentation/p:sldIdLst',
                       namespaces=pptXml.nsmap)
    if len(ret) != 1:
        raise RuntimeError('Broken PPTX file?!')
    else:
        sldRoot = ret[0]

    # Get slide id list
    ret = pptx_slide_get_id_list(sldRoot)
    for i in range(len(ret)):
        sldIdMap[ret[i]] = i + 1

    return sldIdMap


def pptx_get_outline(pptPath):

    # Get presentation title
    prs = Presentation(pptPath)
    if len(prs.slides) == 0:
        raise RuntimeError('Empty PPTX file?!')
    title = prs.slides[0].shapes.title.text

    # Open persentatin xml
    pptXml = etree.fromstring(ZipFile(pptPath).read('ppt/presentation.xml'))

    # Get map of slide id and index
    sldIdMap = pptx_slide_get_id_index_map(pptXml)

    # Get slide information tree
    sldList = pptx_slide_get_list(pptXml)

    # Generate outline
    outline = {
        'title': title,
        'content': [],
        'child': []
    }

    for sect in sldList:
        sectTitle = sect['name']
        sectContent = [sldIdMap[sldId] for sldId in sect['list']]

        outline['child'].append({
            'title': sectTitle,
            'content': sectContent,
            'child': []
        })

    return outline


def _pptx_print_outline(root, level):
    for node in root:
        if node['title'] != None:
            for i in range(level):
                sys.stdout.write('  ')
            print('===', node['title'], '===')
        for content in node['content']:
            for i in range(level):
                sys.stdout.write('  ')
            print(content)
        print()

        _pptx_print_outline(node['child'], level + 1)


def pptx_print_outline(outline):
    _pptx_print_outline([outline], 0)


def pptx_slide_export(pptPath, imgDir, format):

    # Open PowerPoint application
    ppt = Dispatch('PowerPoint.Application')
    ppt.Visible = 1

    # Open presentation session
    pptSess = ppt.Presentations.Open(realpath(pptPath))

    # Export slides as images
    pptSess.Export(realpath(imgDir), format)

    # Close application
    ppt.Quit()


# Test
if __name__ == '__main__':

    import uuid

    fPath = './pptx_file/test.pptx'

    outline = pptx_get_outline(fPath)
    pptx_print_outline(outline)

    pptx_slide_export(fPath, './pptx_file/' + str(uuid.uuid4()), 'PNG')
